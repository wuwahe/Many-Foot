"""
pptx_helpers.py — Reusable drawing primitives for python-pptx scripts.

Install to working directory:
    cp /workspace/skills/create-pptx/scripts/pptx_helpers.py /workspace/code/pptx_helpers.py

Then import:
    import sys; sys.path.insert(0, '/workspace/code')
    from pptx_helpers import bg_rect, h_line, v_line, textbox, oval_shape, \
                              rounded_rect, diag_connector, add_fade_transition, \
                              multiline_textbox, stat_card, icon_card

All position/size arguments are in EMU (English Metric Units).
  1 pt  = 12700 EMU
  1 cm  ≈ 360000 EMU
  16:9  = 12192000 × 6858000 EMU
"""

from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn


# ── Internal helpers ──────────────────────────────────────────────────────────

def _max_existing_id(slide) -> int:
    """Return the highest shape id currently on the slide."""
    return max(
        (int(el.get('id')) for el in slide.element.iter()
         if el.get('id') and el.get('id').isdigit()),
        default=1,
    )


# ── Background ────────────────────────────────────────────────────────────────

def bg_rect(slide, color: RGBColor, sw: int, sh: int):
    """Fill the entire slide with a solid color."""
    shp = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(sw), Emu(sh))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


# ── Lines ─────────────────────────────────────────────────────────────────────

def h_line(slide, x1: int, y: int, x2: int, color: RGBColor, width_pt: float = 1.5):
    """Draw a horizontal line as a thin filled rectangle."""
    h = max(int(width_pt * 12700), 6000)
    shp = slide.shapes.add_shape(1, Emu(x1), Emu(y - h // 2), Emu(x2 - x1), Emu(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def v_line(slide, x: int, y1: int, y2: int, color: RGBColor, width_pt: float = 0.8):
    """Draw a vertical line as a thin filled rectangle."""
    w = max(int(width_pt * 12700), 4000)
    shp = slide.shapes.add_shape(1, Emu(x - w // 2), Emu(y1), Emu(w), Emu(y2 - y1))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def diag_connector(slide, x1: int, y1: int, x2: int, y2: int,
                   color_rgb: tuple, width_pt: float = 0.6) -> int:
    """
    Add a diagonal connector line using raw XML.
    color_rgb: (r, g, b) integers 0-255.
    Returns the shape id of the connector.
    """
    cid = _max_existing_id(slide) + 1
    r, g, b = color_rgb
    dx, dy = abs(x2 - x1) or 1, abs(y2 - y1) or 1
    lx, ly = min(x1, x2), min(y1, y2)
    flip_h = 'flipH="1"' if (x2 < x1) != (y2 < y1) else ''
    xml = (
        f'<p:cxnSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        f' xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
        f'<p:nvCxnSpPr>'
        f'<p:cNvPr id="{cid}" name="conn{cid}"/>'
        f'<p:cNvCxnSpPr/><p:nvPr/>'
        f'</p:nvCxnSpPr>'
        f'<p:spPr>'
        f'<a:xfrm {flip_h}><a:off x="{lx}" y="{ly}"/><a:ext cx="{dx}" cy="{dy}"/></a:xfrm>'
        f'<a:prstGeom prst="line"><a:avLst/></a:prstGeom>'
        f'<a:ln w="{max(int(width_pt * 12700), 6350)}">'
        f'<a:solidFill><a:srgbClr val="{r:02X}{g:02X}{b:02X}"/></a:solidFill>'
        f'<a:round/></a:ln>'
        f'</p:spPr>'
        f'</p:cxnSp>'
    )
    slide.shapes._spTree.append(parse_xml(xml.encode()))
    return cid


# ── Text ──────────────────────────────────────────────────────────────────────

def textbox(slide, x: int, y: int, w: int, h: int,
            text: str, size: float, color: RGBColor,
            bold: bool = False, align=PP_ALIGN.LEFT):
    """Add a single-line textbox (word_wrap=False)."""
    txb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = txb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return txb


def multiline_textbox(slide, x: int, y: int, w: int, h: int,
                      lines: list, size: float, color: RGBColor,
                      bold: bool = False, align=PP_ALIGN.LEFT,
                      line_spacing_pt: float = 0,
                      bg_color: RGBColor = None,
                      padding_emu: int = 55000):
    """
    Add a textbox with multiple lines (one paragraph per item in `lines`).
    lines: list of str, or list of (str, RGBColor) tuples for per-line color.
    bg_color: optional background fill for the box.
    """
    if bg_color:
        shp = slide.shapes.add_shape(1, Emu(x), Emu(y), Emu(w), Emu(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = bg_color
        shp.line.fill.background()

    txb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = txb.text_frame
    tf.word_wrap = True
    if bg_color:
        tf.margin_left = tf.margin_right = Emu(padding_emu)
        tf.margin_top  = tf.margin_bottom = Emu(padding_emu // 2)

    for i, line in enumerate(lines):
        # Support (text, color) tuples for per-line coloring
        if isinstance(line, tuple):
            text, line_color = line
        else:
            text, line_color = line, color

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing_pt:
            p.space_before = Pt(line_spacing_pt)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = line_color
        run.font.bold = bold
    return txb


# ── Shapes ────────────────────────────────────────────────────────────────────

def oval_shape(slide, cx: int, cy: int, size: int, color: RGBColor) -> int:
    """Add a filled circle centered at (cx, cy). Returns shape_id."""
    shp = slide.shapes.add_shape(9, Emu(cx - size // 2), Emu(cy - size // 2),
                                 Emu(size), Emu(size))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp.shape_id


def rounded_rect(slide, x: int, y: int, w: int, h: int,
                 fill: RGBColor, border: RGBColor = None,
                 border_pt: float = 0.0, radius_emu: int = 60000):
    """
    Add a rounded rectangle. Returns the shape.
    Set border=None and border_pt=0 for no border (common on dark backgrounds).
    radius_emu: corner radius as EMU value (default 60000 ≈ subtle rounding).
    """
    shp = slide.shapes.add_shape(5, Emu(x), Emu(y), Emu(w), Emu(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border and border_pt > 0:
        shp.line.color.rgb = border
        shp.line.width = Pt(border_pt)
    else:
        shp.line.fill.background()
    # Adjust corner radius via XML
    prstGeom = shp.element.spPr.find(qn('a:prstGeom'))
    if prstGeom is not None:
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is not None:
            gd = parse_xml(
                f'<a:gd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
                f' name="adj" fmla="val {radius_emu}"/>'
            )
            avLst.append(gd)
    return shp


def stat_card(slide, x: int, y: int, w: int, h: int,
              value: str, label: str,
              accent: RGBColor, bg: RGBColor,
              value_size: float = 32, label_size: float = 11):
    """
    Draw a KPI / stat card: big number on top, label below.
    Common in data summary slides.
    """
    rounded_rect(slide, x, y, w, h, fill=bg, radius_emu=80000)

    # Accent top bar
    bar_h = max(h // 12, 40000)
    rounded_rect(slide, x, y, w, bar_h, fill=accent, radius_emu=80000)
    # Flat bottom of bar (overdraw)
    slide.shapes.add_shape(1, Emu(x), Emu(y + bar_h // 2),
                            Emu(w), Emu(bar_h // 2)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = accent
    slide.shapes[-1].line.fill.background()

    # Value text
    textbox(slide, x, y + h // 5, w, h // 2,
            value, value_size, RGBColor(0xFF, 0xFF, 0xFF),
            bold=True, align=PP_ALIGN.CENTER)

    # Label text
    textbox(slide, x, y + h * 2 // 3, w, h // 4,
            label, label_size, RGBColor(0xAA, 0xBB, 0xCC),
            align=PP_ALIGN.CENTER)


def icon_card(slide, x: int, y: int, w: int, h: int,
              icon_char: str, title: str, body: str,
              accent: RGBColor, bg: RGBColor,
              icon_size: float = 28, title_size: float = 13,
              body_size: float = 10):
    """
    Draw an info card with an icon character, title, and body text.
    icon_char: a Unicode char or emoji used as icon (e.g. '✦', '→', '📊').
    """
    rounded_rect(slide, x, y, w, h, fill=bg, radius_emu=80000)

    icon_h = h // 4
    # Icon circle
    oval_shape(slide, x + w // 2, y + icon_h // 2 + 60000,
               min(icon_h, 300000), accent)
    # Icon char
    textbox(slide, x, y + 30000, w, icon_h,
            icon_char, icon_size, RGBColor(0xFF, 0xFF, 0xFF),
            bold=True, align=PP_ALIGN.CENTER)
    # Title
    textbox(slide, x + 40000, y + icon_h + 30000, w - 80000, h // 5,
            title, title_size, RGBColor(0xFF, 0xFF, 0xFF),
            bold=True, align=PP_ALIGN.CENTER)
    # Body
    multiline_textbox(slide, x + 40000, y + icon_h + h // 5 + 30000,
                      w - 80000, h - icon_h - h // 5 - 60000,
                      body.split('\n'), body_size,
                      RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.CENTER)


# ── Transitions ───────────────────────────────────────────────────────────────

def add_fade_transition(slide, spd: str = 'med'):
    """
    Inject a cross-fade slide transition.
    Works in both Microsoft PowerPoint and WPS.
    spd: 'slow' | 'med' | 'fast'
    Call this on every slide before saving.
    """
    PNS = 'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
    ANS = 'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
    xml = (
        f'<p:transition {PNS} {ANS} spd="{spd}" advClick="1">'
        f'<p:fade/>'
        f'</p:transition>'
    )
    elem = parse_xml(xml.encode())
    sld = slide.element
    old = sld.find(qn('p:transition'))
    if old is not None:
        sld.remove(old)
    cSld = sld.find(qn('p:cSld'))
    cSld.addnext(elem)
